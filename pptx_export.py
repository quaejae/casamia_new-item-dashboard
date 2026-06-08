# -*- coding: utf-8 -*-
"""PPTX export — DashboardTable(들) → A4 가로 PPTX.

담당자 1명 = 슬라이드 1장. 예시 양식과 동일한 구조:
  주제목(텍스트) → 소제목(➤담당자N) → 24열 표(헤더2행 + 데이터 + 요약행)
표는 일반 표 객체로 출력되어 팀장이 PowerPoint에서 자유롭게 수정 가능.
'이미지' 열에는 셀 위치에 맞춰 Picture 도형을 겹쳐 배치(없으면 빈 칸).
"""
from __future__ import annotations

import io
from typing import List, Optional

from pptx import Presentation
from pptx.util import Mm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml

import config
from schema import DashboardTable
from transform import launch_arrow_span, current_month_index

# 표 열 구성: 정보 10열 + 13개월 + 비고 = 24
N_INFO = 10
N_MONTH = 13
N_COLS = N_INFO + N_MONTH + 1   # 24
MONTH_BASE = N_INFO            # 월 컬럼 시작 인덱스(10)
NOTE_COL = N_COLS - 1          # 비고(23)

INFO_HEADERS = ["NO", "시리즈명", "이미지", "소재", "목표\n판매가",
                "목표\n매출", "컬러", "SKU수", "제조처", "담당자"]


def _hex(rgb: str) -> RGBColor:
    return RGBColor.from_string(rgb)


def _set_cell(cell, text="", *, size=None, bold=True, fill=None,
              align=PP_ALIGN.CENTER, color="000000"):
    """셀 텍스트·서식 설정."""
    size = size or config.FONT_BODY
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Pt(1); cell.margin_right = Pt(1)
    cell.margin_top = Pt(1); cell.margin_bottom = Pt(1)
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = _hex(fill)
    else:
        cell.fill.background()
    tf = cell.text_frame
    tf.word_wrap = True
    # 여러 줄 텍스트 처리
    lines = str(text).split("\n")
    tf.clear()
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = config.FONT_NAME
        run.font.color.rgb = _hex(color)


def _apply_gradient(shape, c_from: str, c_to: str):
    """도형에 가로 그라데이션(좌:c_from → 우:c_to) 적용, 테두리 제거.

    스키마상 fill 은 geometry 다음, ln 앞에 와야 PowerPoint 가 인식한다.
    또한 <p:style> 의 fillRef(테마색) 가 우선되지 않도록 style 의 fill 참조를 제거.
    """
    spPr = shape._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill",
                "a:pattFill", "a:grpFill"):
        for el in spPr.findall(__qn(tag)):
            spPr.remove(el)
    shape.line.fill.background()   # 먼저 a:ln(noFill) 을 올바른 위치에 생성
    grad = parse_xml(
        '<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:gsLst>'
        f'<a:gs pos="0"><a:srgbClr val="{c_from}"/></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{c_to}"/></a:gs>'
        '</a:gsLst><a:lin ang="0" scaled="0"/></a:gradFill>'
    )
    # geometry(prstGeom/custGeom) 바로 뒤에 삽입
    geom = spPr.find(__qn("a:prstGeom"))
    if geom is None:
        geom = spPr.find(__qn("a:custGeom"))
    if geom is not None:
        geom.addnext(grad)
    else:
        spPr.insert(0, grad)


def __qn(tag: str):
    from pptx.oxml.ns import qn
    return qn(tag)


def _col_widths(usable_emu: int) -> List[int]:
    total = sum(config.COL_WEIGHTS)
    return [int(usable_emu * w / total) for w in config.COL_WEIGHTS]


def _milestone_fill(text: str) -> Optional[str]:
    for kw, color in config.MILESTONE_COLORS.items():
        if kw in text:
            return color
    return None


def _text_color(text: str) -> str:
    """발주/출시 마일스톤은 주황, 그 외 검정."""
    t = (text or "").strip()
    if t in config.ORANGE_TEXT_EXACT:
        return config.ORANGE_TEXT_COLOR
    if any(kw in t for kw in config.ORANGE_TEXT_CONTAINS):
        return config.ORANGE_TEXT_COLOR
    return "000000"


def _set_cell_border(cell):
    """셀 4면에 회색 테두리 적용 (tcPr 맨 앞에 lnL/R/T/B 삽입)."""
    w = int(config.BORDER_WIDTH_PT * 12700)
    color = config.BORDER_COLOR
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for el in tcPr.findall(__qn(tag)):
            tcPr.remove(el)
    # 순서상 lnL,lnR,lnT,lnB 는 tcPr 의 가장 앞에 와야 함 → 역순 insert(0)
    for tag in ("a:lnB", "a:lnT", "a:lnR", "a:lnL"):
        ln = parse_xml(
            f'<{tag} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            f'w="{w}" cap="flat" cmpd="sng" algn="ctr">'
            f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
            f'<a:prstDash val="solid"/></{tag}>'
        )
        tcPr.insert(0, ln)


def _add_title_block(slide, main_title, subtitle_text, left_emu, top_emu, width_emu):
    """주제목(고정 위치) → 장식선 → 소제목(표 바로 위). 표 시작 top 반환."""
    # 주제목 (좌상단 고정)
    tb = slide.shapes.add_textbox(left_emu, top_emu, width_emu, Mm(8))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = main_title
    r.font.size = Pt(config.FONT_TITLE); r.font.bold = True
    r.font.name = config.FONT_NAME; r.font.color.rgb = _hex("000000")
    # 주제목 아래 장식선
    if config.TITLE_RULE_ENABLED:
        ln = slide.shapes.add_connector(2, left_emu, top_emu + Mm(8),
                                        left_emu + width_emu, top_emu + Mm(8))
        ln.line.color.rgb = _hex(config.TITLE_RULE_COLOR)
        ln.line.width = Pt(config.TITLE_RULE_WIDTH_PT)
        ln.shadow.inherit = False
    # 소제목 (표 바로 위)
    sb = slide.shapes.add_textbox(left_emu, top_emu + Mm(9), width_emu, Mm(7))
    p2 = sb.text_frame.paragraphs[0]
    r2 = p2.add_run(); r2.text = subtitle_text
    r2.font.size = Pt(config.FONT_SUBTITLE); r2.font.bold = True
    r2.font.name = config.FONT_NAME; r2.font.color.rgb = _hex(config.SUBTITLE_COLOR)
    return top_emu + Mm(17)   # 표 시작 top


def _add_footer(slide, left_emu, width_emu, slide_height):
    """하단 장식선 + 'SHINSEGAE CASA' 브랜드 텍스트(슬라이드마다 고정)."""
    if not config.FOOTER_ENABLED:
        return
    y = slide_height - Mm(config.FOOTER_LINE_MARGIN_MM)
    ln = slide.shapes.add_connector(2, left_emu, y, left_emu + width_emu, y)
    ln.line.color.rgb = _hex(config.TITLE_RULE_COLOR)
    ln.line.width = Pt(config.TITLE_RULE_WIDTH_PT)
    ln.shadow.inherit = False
    tb = slide.shapes.add_textbox(left_emu, y, width_emu,
                                  Mm(config.FOOTER_LINE_MARGIN_MM))
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = config.FOOTER_TEXT
    r.font.size = Pt(config.FOOTER_FONT); r.font.bold = True
    r.font.name = config.FONT_NAME; r.font.color.rgb = _hex(config.FOOTER_COLOR)


def _build_header(tbl, show_q: bool):
    # 정보 10열: row0:row1 세로 병합
    for ci, label in enumerate(INFO_HEADERS):
        txt = label
        if ci == 4 and show_q:   # 목표판매가 + (Q)
            txt = label + "\n" + config.Q_SIZE_LABEL
        a = tbl.cell(0, ci); b = tbl.cell(1, ci)
        a.merge(b)
        _set_cell(a, txt, size=config.FONT_HEADER, fill=config.HEADER_FILL)
    # 비고: 세로 병합
    a = tbl.cell(0, NOTE_COL); b = tbl.cell(1, NOTE_COL)
    a.merge(b)
    _set_cell(a, "비고", size=config.FONT_HEADER, fill=config.HEADER_FILL)
    # 분기 헤더(row0) + 월(row1)
    for qlabel, start, length in config.QUARTERS:
        c0 = tbl.cell(0, MONTH_BASE + start)
        if length > 1:
            c1 = tbl.cell(0, MONTH_BASE + start + length - 1)
            c0.merge(c1)
        _set_cell(c0, qlabel, size=config.FONT_HEADER, fill=config.HEADER_FILL)
    for mi, (_, month) in enumerate(config.TIMELINE_MONTHS):
        _set_cell(tbl.cell(1, MONTH_BASE + mi), str(month),
                  size=config.FONT_HEADER, fill=config.HEADER_FILL)


def _build_table(slide, rows, summary, show_q, left_emu, top_emu, usable_emu):
    """행 묶음(rows)을 한 슬라이드의 표로 렌더. summary 가 None 이면 요약행 생략."""
    n_data = len(rows)
    has_summary = summary is not None
    n_rows = 2 + n_data + (1 if has_summary else 0)
    widths = _col_widths(usable_emu)

    gf = slide.shapes.add_table(n_rows, N_COLS, left_emu, top_emu,
                                usable_emu, Mm(10))
    tbl = gf.table
    tbl.first_row = False; tbl.horz_banding = False
    for ci, w in enumerate(widths):
        tbl.columns[ci].width = Emu(w)
    # 행 높이
    tbl.rows[0].height = Pt(config.ROW_H_HEADER)
    tbl.rows[1].height = Pt(config.ROW_H_MONTH)
    for ri in range(2, 2 + n_data):
        tbl.rows[ri].height = Pt(config.ROW_H_DATA)
    if has_summary:
        tbl.rows[n_rows - 1].height = Pt(config.ROW_H_SUMMARY)

    _build_header(tbl, show_q)

    # 데이터 행
    image_jobs = []   # (table_row, bytes, ext) — 표 생성 후 오버레이
    arrow_jobs = []   # (di, source_idx, launch_idx) — 출시 화살표
    for di, row in enumerate(rows):
        ri = 2 + di
        _set_cell(tbl.cell(ri, 0), row.no_label, size=config.FONT_BODY)
        _set_cell(tbl.cell(ri, 1), row.series_cell, size=config.FONT_BODY)
        _set_cell(tbl.cell(ri, 2), "", size=config.FONT_BODY)   # 이미지 칸
        _set_cell(tbl.cell(ri, 3), row.material, size=config.FONT_BODY)
        _set_cell(tbl.cell(ri, 4), row.target_price, size=config.FONT_BODY)
        _set_cell(tbl.cell(ri, 5), row.target_revenue, size=config.FONT_BODY)
        _set_cell(tbl.cell(ri, 6), row.color_count, size=config.FONT_BODY)
        _set_cell(tbl.cell(ri, 7), row.sku_count, size=config.FONT_BODY)
        _set_cell(tbl.cell(ri, 8), row.maker, size=config.FONT_BODY)
        _set_cell(tbl.cell(ri, 9), row.owner, size=config.FONT_BODY)
        for mi in range(N_MONTH):
            txt = row.timeline[mi] if mi < len(row.timeline) else ""
            _set_cell(tbl.cell(ri, MONTH_BASE + mi), txt,
                      size=config.FONT_BODY, fill=_milestone_fill(txt),
                      color=_text_color(txt))
        _set_cell(tbl.cell(ri, NOTE_COL), row.note, size=config.FONT_BODY)
        if row.image:
            image_jobs.append((ri, row.image, row.image_ext))
        if config.ARROW_ENABLED:
            span = launch_arrow_span(row.timeline)
            if span:
                arrow_jobs.append((di, span[0], span[1]))

    # NO 세로 병합 (연속된 빈 no_label 을 위 행과 병합)
    start = 2
    for di in range(1, n_data + 1):
        ri = 2 + di
        is_new = (di == n_data) or (rows[di].no_label != "")
        if is_new:
            if ri - 1 > start:
                tbl.cell(start, 0).merge(tbl.cell(ri - 1, 0))
            start = ri

    # 요약행 (마지막 페이지에만)
    if has_summary:
        sr = n_rows - 1
        s = summary
        a = tbl.cell(sr, 0); b = tbl.cell(sr, N_INFO - 1)
        a.merge(b)
        _set_cell(a, s.label, size=config.FONT_SUMMARY, fill=config.SUMMARY_FILL,
                  align=PP_ALIGN.CENTER)
        for qi, (_, start_m, length) in enumerate(config.QUARTERS):
            c0 = tbl.cell(sr, MONTH_BASE + start_m)
            if length > 1:
                c1 = tbl.cell(sr, MONTH_BASE + start_m + length - 1)
                c0.merge(c1)
            _set_cell(c0, s.per_quarter[qi],
                      size=config.FONT_SUMMARY, fill=config.SUMMARY_FILL)
        _set_cell(tbl.cell(sr, NOTE_COL), f"{s.series_count}({s.total_sku})",
                  size=config.FONT_SUMMARY, fill=config.SUMMARY_FILL)

    # 모든 셀에 회색 테두리 적용
    for ri in range(n_rows):
        for ci in range(N_COLS):
            _set_cell_border(tbl.cell(ri, ci))

    return gf, widths, image_jobs, arrow_jobs


def _add_month_highlight(slide, table_left, table_top, widths, n_data, month_idx):
    """오늘 날짜의 월 열 데이터 영역에 네이비 사각 테두리 강조(채움 없음)."""
    cum = [table_left]
    for w in widths:
        cum.append(cum[-1] + w)
    header_h = Pt(config.ROW_H_HEADER) + Pt(config.ROW_H_MONTH)
    box_left = cum[MONTH_BASE + month_idx]
    box_top = table_top + header_h
    box_w = widths[MONTH_BASE + month_idx]
    box_h = int(Pt(config.ROW_H_DATA) * n_data)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 int(box_left), int(box_top), int(box_w), box_h)
    box.fill.background()
    box.line.color.rgb = _hex(config.CURRENT_MONTH_COLOR)
    box.line.width = Pt(config.CURRENT_MONTH_WIDTH_PT)
    box.shadow.inherit = False
    return box


def _add_launch_arrows(slide, table_left, table_top, widths, arrow_jobs, table_gf):
    """진행기간 화살표(기획→출시 직전)를 표 '뒤'에 배치.

    셀 배경이 투명하므로 화살표가 셀을 가로질러 비쳐, 프로젝트 진행기간이 보인다.
    """
    h_header = Pt(config.ROW_H_HEADER) + Pt(config.ROW_H_MONTH)
    h_data = Pt(config.ROW_H_DATA)
    arrow_h = int(h_data * config.ARROW_HEIGHT_RATIO)
    cum = [table_left]
    for w in widths:
        cum.append(cum[-1] + w)
    spTree = slide.shapes._spTree
    table_el = table_gf._element
    for di, start, end in arrow_jobs:
        x_start = cum[MONTH_BASE + start]
        x_end = cum[MONTH_BASE + end + 1]   # 출시 직전 칸의 오른쪽 끝
        width = x_end - x_start
        if width <= 0:
            continue
        row_top = table_top + h_header + h_data * di
        y = int(row_top + (h_data - arrow_h) / 2)
        shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                     int(x_start), y, int(width), arrow_h)
        shp.shadow.inherit = False
        _apply_gradient(shp, config.ARROW_COLOR_FROM, config.ARROW_COLOR_TO)
        # z-order: 표 graphicFrame 바로 앞(=뒤쪽)으로 이동
        spTree.remove(shp._element)
        table_el.addprevious(shp._element)


def _overlay_images(slide, table_left, table_top, widths, image_jobs):
    """이미지 칸(col 2) 위에 Picture 를 셀 크기에 맞춰 배치."""
    col2_left = table_left + sum(widths[:2])
    col2_w = widths[2]
    # 행 높이 누적: 헤더2행 높이 + 데이터행 높이
    h_header = Pt(config.ROW_H_HEADER) + Pt(config.ROW_H_MONTH)
    h_data = Pt(config.ROW_H_DATA)
    pad = Pt(2)
    for ri, data, ext in image_jobs:
        di = ri - 2
        cell_top = table_top + h_header + h_data * di
        box_left = col2_left + pad
        box_top = cell_top + pad
        box_w = col2_w - 2 * pad
        box_h = h_data - 2 * pad
        stream = io.BytesIO(data)
        # 가로 기준으로 넣고 비율 유지(세로가 넘치면 세로 기준)
        pic = slide.shapes.add_picture(stream, box_left, box_top, width=box_w)
        if pic.height > box_h:
            scale = box_h / pic.height
            pic.height = int(pic.height * scale)
            pic.width = int(pic.width * scale)
        # 셀 중앙 정렬
        pic.left = int(col2_left + (col2_w - pic.width) / 2)
        pic.top = int(cell_top + (h_data - pic.height) / 2)


def _paginate_rows(rows, capacity):
    """행을 시리즈 블록(no_label!='' 시작) 단위로 페이지에 채워 분할.

    시리즈가 페이지 경계에서 쪼개지지 않도록 블록째 배치.
    """
    blocks = []
    cur = []
    for r in rows:
        if r.no_label != "" and cur:
            blocks.append(cur); cur = []
        cur.append(r)
    if cur:
        blocks.append(cur)
    pages, page = [], []
    for blk in blocks:
        if page and len(page) + len(blk) > capacity:
            pages.append(page); page = []
        page.extend(blk)
    if page:
        pages.append(page)
    return pages or [[]]


def _build_presentation(tables: List[DashboardTable]) -> Presentation:
    prs = Presentation()
    prs.slide_width = Mm(config.A4_LANDSCAPE_WIDTH_MM)
    prs.slide_height = Mm(config.A4_LANDSCAPE_HEIGHT_MM)
    blank = prs.slide_layouts[6]   # 빈 레이아웃

    margin = Mm(config.PAGE_MARGIN_MM)
    usable_w = prs.slide_width - 2 * margin
    table_top = margin + Mm(17)
    bottom_line_y = prs.slide_height - Mm(config.FOOTER_LINE_MARGIN_MM)

    # 한 슬라이드(장식선 사이)에 들어갈 데이터 행 수 = 페이지 용량
    avail = bottom_line_y - table_top - Mm(config.PAGE_GAP_MM)
    header_h = Pt(config.ROW_H_HEADER) + Pt(config.ROW_H_MONTH)
    capacity = max(1, int((avail - header_h - Pt(config.ROW_H_SUMMARY))
                          / Pt(config.ROW_H_DATA)))

    cur_mi = current_month_index() if config.CURRENT_MONTH_HIGHLIGHT else None

    for table in tables:
        pages = _paginate_rows(table.rows, capacity)
        total = len(pages)
        for pi, page_rows in enumerate(pages):
            slide = prs.slides.add_slide(blank)
            subtitle = f"{config.SUBTITLE_PREFIX}{table.owner}"
            if total > 1:
                subtitle += f" ({pi + 1}/{total})"
            _add_title_block(slide, table.main_title, subtitle,
                             margin, margin, usable_w)
            summary = table.summary if pi == total - 1 else None
            gf, widths, image_jobs, arrow_jobs = _build_table(
                slide, page_rows, summary, table.show_q_label,
                margin, table_top, usable_w)
            if arrow_jobs:
                _add_launch_arrows(slide, margin, table_top, widths, arrow_jobs, gf)
            if image_jobs:
                _overlay_images(slide, margin, table_top, widths, image_jobs)
            if cur_mi is not None and page_rows:
                _add_month_highlight(slide, margin, table_top, widths,
                                     len(page_rows), cur_mi)
            _add_footer(slide, margin, usable_w, prs.slide_height)
    return prs


def build_pptx(tables: List[DashboardTable], out_path: str) -> str:
    _build_presentation(tables).save(out_path)
    return out_path


def build_pptx_bytes(tables: List[DashboardTable]) -> bytes:
    """PPTX 를 메모리(bytes)로 반환 — Streamlit 다운로드 버튼용."""
    buf = io.BytesIO()
    _build_presentation(tables).save(buf)
    return buf.getvalue()


if __name__ == "__main__":
    import sys
    import data_loader, transform
    src = sys.argv[1] if len(sys.argv) > 1 else None
    out = sys.argv[2] if len(sys.argv) > 2 else "대시보드_export.pptx"
    tables = transform.build_all(data_loader.load(src))
    path = build_pptx(tables, out)
    print("생성 완료:", path, "| 슬라이드", len(tables), "장")
